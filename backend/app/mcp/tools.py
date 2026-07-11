"""MCP tool catalog (spec 05 §2) — real extractors. Each returns structured content
or a structured error (never a silent empty). Binaries are stored via the BlobStore
and returned as blob refs. Extractors are PURE: given a ref/url they return content;
they don't decide relevance (that's the Scout/Auditor).

Priority set built first (D20): web search/scrape, arXiv paper download, PDF text,
vision figures, slides. The rest degrade to best-effort.
"""
from __future__ import annotations

import io
import re
from typing import Any

import httpx

from app.blobstore import get_blobstore

_UA = "Mozilla/5.0 (compatible; Seekhai_test/1.0; +local research)"
_TIMEOUT = 20.0


def _err(tool: str, msg: str, **extra) -> dict:
    return {"error": True, "tool": tool, "message": msg, **extra}


def _extractor_shape(source_id: str, mime: str, **kw) -> dict[str, Any]:
    base = {"source_id": source_id, "mime": mime, "text_chunks": [], "figures": [],
            "tables": [], "citations": [], "meta": {}}
    base.update(kw)
    return base


# --- discovery ---------------------------------------------------------------
# stock-image / unrelated domains that collide on generic terms like "model"
_JUNK = ("sketchfab.com", "unsplash.com", "models.com", "shutterstock", "istockphoto",
         "modelbank", "pinterest", "gettyimages", "freepik", "123rf")


def _wikipedia_search(query: str, max_results: int = 3) -> list[dict]:
    """Reliable keyless search via the Wikipedia REST/opensearch API."""
    try:
        with httpx.Client(timeout=_TIMEOUT, headers={"User-Agent": _UA}) as c:
            r = c.get("https://en.wikipedia.org/w/api.php",
                      params={"action": "opensearch", "search": query,
                              "limit": max_results, "namespace": 0, "format": "json"})
            r.raise_for_status()
            data = r.json()
        titles, _descs, urls = data[1], data[2], data[3]
        return [{"title": t, "url": u, "snippet": None, "source": "wikipedia"}
                for t, u in zip(titles, urls)]
    except Exception:
        return []


def _ddgs_search(query: str, max_results: int = 5) -> list[dict]:
    try:
        from ddgs import DDGS

        out = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                url = r.get("href") or ""
                if any(j in url.lower() for j in _JUNK):
                    continue
                out.append({"title": r.get("title"), "url": url, "snippet": r.get("body")})
        return out
    except Exception:
        return []


def web_search(query: str, max_results: int = 6, since: str | None = None,
               course_id: str | None = None) -> dict:
    """Live web search. Primary: OpenRouter's built-in web plugin (Exa-backed) — high
    quality, on-topic. Falls back to DuckDuckGo (junk-filtered) + Wikipedia if that
    provider errors, so scouting still works offline-ish."""
    from app import events

    # primary: real search provider
    try:
        from app.llm import web_search as or_search

        out = or_search(query, max_results=max_results, course_id=course_id)
        results = [r for r in out["results"] if r.get("url")
                   and not any(j in r["url"].lower() for j in _JUNK)]
        if results:
            events.emit(course_id, "scouting", "web_search",
                        f'web_search (OpenRouter/Exa) "{query[:60]}" → {len(results)} results',
                        {"provider": "openrouter_web", "results": len(results),
                         "cost": round(out.get("cost", 0.0), 5)})
            return {"results": results[:max_results]}
    except Exception as e:
        events.emit(course_id, "scouting", "warn",
                    f"primary web_search failed ({type(e).__name__}); falling back to DDG+Wikipedia")

    # fallback: DDGS + Wikipedia
    results, seen = [], set()
    for r in _ddgs_search(query, max_results) + _wikipedia_search(query, 3):
        u = r.get("url")
        if u and u not in seen:
            seen.add(u)
            results.append(r)
    events.emit(course_id, "scouting", "web_search",
                f'web_search (fallback DDG+Wiki) "{query[:60]}" → {len(results)} results')
    return {"results": results[:max_results]}


def paper_search(query: str, since: str | None = None, max_results: int = 5) -> dict:
    """Search arXiv, date-filterable for 'recent' (spec 05 §5)."""
    try:
        import arxiv

        client = arxiv.Client(page_size=max_results, delay_seconds=1, num_retries=2)
        search = arxiv.Search(query=query, max_results=max_results,
                              sort_by=arxiv.SortCriterion.SubmittedDate)
        results = []
        for r in client.results(search):
            published = r.published.date().isoformat() if r.published else None
            if since and published and published < since:
                continue
            results.append({
                "title": r.title, "url": r.pdf_url, "entry_id": r.entry_id,
                "published": published, "authors": [a.name for a in r.authors[:6]],
                "summary": (r.summary or "").strip()[:1500], "type": "paper",
            })
        return {"results": results}
    except Exception as e:
        return _err("paper_search", f"{type(e).__name__}: {e}", results=[])


# --- fetch -------------------------------------------------------------------
def file_fetcher(url: str) -> dict:
    """Download an arbitrary document → blob ref + detected MIME."""
    try:
        with httpx.Client(follow_redirects=True, timeout=_TIMEOUT,
                          headers={"User-Agent": _UA}) as c:
            resp = c.get(url)
            resp.raise_for_status()
            mime = resp.headers.get("content-type", "application/octet-stream").split(";")[0]
            blob_id = get_blobstore().put("source", mime, resp.content)
            return {"blob_ref": blob_id, "mime": mime, "byte_len": len(resp.content), "url": url}
    except Exception as e:
        return _err("file_fetcher", f"{type(e).__name__}: {e}", url=url)


def paper_downloader(identifier: str) -> dict:
    """Download a paper (arXiv id / url) → PDF blob ref + metadata."""
    url = identifier
    if not url.lower().endswith(".pdf") and "arxiv.org/abs/" in url:
        url = url.replace("/abs/", "/pdf/") + ".pdf"
    f = file_fetcher(url)
    if f.get("error"):
        return f
    return {"pdf_ref": f["blob_ref"], "mime": f["mime"], "url": url, "byte_len": f["byte_len"]}


# --- format-specific extractors ---------------------------------------------
def _clean_chunks(text: str, source_id: str, max_chunks: int = 24) -> list[dict]:
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if len(p.strip()) > 120]
    return [{"id": f"{source_id[:8]}-{i}", "text": p[:2000], "ordinal": i}
            for i, p in enumerate(paras[:max_chunks])]


def web_scrape(url: str) -> dict:
    """Fetch + clean a page's main prose (trafilatura, bs4 fallback)."""
    try:
        with httpx.Client(follow_redirects=True, timeout=_TIMEOUT,
                          headers={"User-Agent": _UA}) as c:
            resp = c.get(url)
            resp.raise_for_status()
            html = resp.text
        text, title, published = None, None, None
        try:
            import trafilatura

            text = trafilatura.extract(html, include_comments=False, include_tables=True)
            md = trafilatura.extract_metadata(html)
            if md:
                title = md.title
                published = md.date
        except Exception:
            pass
        if not text:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            if soup.title:
                title = soup.title.get_text(strip=True)
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text("\n", strip=True)
        return {"url": url, "title": title, "text": text or "", "published": published}
    except Exception as e:
        return _err("web_scrape", f"{type(e).__name__}: {e}", url=url)


def html_article_extractor(url: str) -> dict:
    s = web_scrape(url)
    if s.get("error"):
        return s
    return _extractor_shape(url, "text/html",
                            text_chunks=_clean_chunks(s.get("text", ""), url),
                            meta={"title": s.get("title"), "published": s.get("published")})


def pdf_extractor(ref: str) -> dict:
    """Extract text (per page) from a PDF blob."""
    try:
        got = get_blobstore().get(ref)
        if got is None:
            return _err("pdf_extractor", "blob not found", source_id=ref)
        _, data = got
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        chunks: list[dict] = []
        for pageno, page in enumerate(reader.pages[:30]):
            txt = (page.extract_text() or "").strip()
            if len(txt) > 100:
                chunks.append({"id": f"{ref[:8]}-p{pageno}", "text": txt[:3000],
                               "page": pageno + 1, "ordinal": pageno})
        meta = {}
        try:
            info = reader.metadata or {}
            meta = {"title": info.get("/Title"), "authors": info.get("/Author")}
        except Exception:
            pass
        return _extractor_shape(ref, "application/pdf", text_chunks=chunks, meta=meta,
                                page_count=len(reader.pages))
    except Exception as e:
        return _err("pdf_extractor", f"{type(e).__name__}: {e}", source_id=ref)


def vision_image_extractor(ref: str, page: int | None = None, source_url: str | None = None) -> dict:
    """Extract figures/diagrams + captions from a PDF page or image using a vision model.
    Renders the source to an image and asks the multimodal model to describe figures."""
    try:
        got = get_blobstore().get(ref)
        if got is None:
            return _err("vision_image_extractor", "blob not found", figures=[])
        mime, data = got
        image_bytes, image_mime = None, mime
        if "pdf" in mime.lower():
            try:
                import pypdfium2 as pdfium  # optional render dependency

                pdf = pdfium.PdfDocument(data)
                pil = pdf[page or 0].render(scale=2).to_pil()
                buf = io.BytesIO(); pil.save(buf, format="PNG")
                image_bytes, image_mime = buf.getvalue(), "image/png"
            except Exception:
                return {"figures": [], "note": "pdf render unavailable (pypdfium2 not installed)"}
        else:
            image_bytes = data
        from app.llm import complete_json, image_part

        parts = [
            {"type": "text", "text": "List the distinct figures/diagrams/charts in this image. "
             "For each, give a short caption and its kind (diagram|chart|photo|schematic). "
             "Return JSON {\"figures\":[{\"caption\":...,\"kind\":...}]}. If none, empty list."},
            image_part(image_bytes, image_mime),
        ]
        data_out, _ = complete_json("vision_image_extractor",
                                    "You extract figures from images. Output only JSON.",
                                    parts, phase="scouting", max_tokens=1000)
        figs = [{"image_ref": ref, "caption": f.get("caption"),
                 "kind": f.get("kind", "diagram"), "page": page, "source_url": source_url}
                for f in (data_out.get("figures") or [])]
        return {"figures": figs}
    except Exception as e:
        return _err("vision_image_extractor", f"{type(e).__name__}: {e}", figures=[])


def slides_extractor(ref: str) -> dict:
    try:
        got = get_blobstore().get(ref)
        if got is None:
            return _err("slides_extractor", "blob not found", source_id=ref)
        _, data = got
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        chunks = []
        for i, slide in enumerate(prs.slides):
            texts = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False)]
            joined = "\n".join(t for t in texts if t.strip())
            if joined.strip():
                chunks.append({"id": f"{ref[:8]}-s{i}", "text": joined[:2000], "ordinal": i})
        return _extractor_shape(ref, "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                text_chunks=chunks)
    except Exception as e:
        return _err("slides_extractor", f"{type(e).__name__}: {e}", source_id=ref)


def doc_extractor(ref: str) -> dict:
    try:
        got = get_blobstore().get(ref)
        if got is None:
            return _err("doc_extractor", "blob not found", source_id=ref)
        _, data = got
        from docx import Document

        doc = Document(io.BytesIO(data))
        text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return _extractor_shape(ref, "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                text_chunks=_clean_chunks(text, ref))
    except Exception as e:
        return _err("doc_extractor", f"{type(e).__name__}: {e}", source_id=ref)


def illustration_scraper(url: str) -> dict:
    """Find clean illustrations/diagrams on an image-heavy page (img tags + alt/caption)."""
    try:
        with httpx.Client(follow_redirects=True, timeout=_TIMEOUT,
                          headers={"User-Agent": _UA}) as c:
            html = c.get(url).text
        from urllib.parse import urljoin

        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        images = []
        for img in soup.find_all("img")[:20]:
            src = img.get("src") or img.get("data-src")
            if not src:
                continue
            src = urljoin(url, src)
            alt = img.get("alt") or ""
            if src.lower().split("?")[0].endswith((".svg", ".png", ".jpg", ".jpeg", ".webp")):
                images.append({"url": src, "alt": alt, "caption": alt, "license_hint": None})
        return {"images": images}
    except Exception as e:
        return _err("illustration_scraper", f"{type(e).__name__}: {e}", images=[])


# best-effort long tail of formats (kept in the catalog, spec 05 §2)
def transcript_extractor(ref: str) -> dict:
    return _extractor_shape(ref, "text/vtt", note="transcript extraction not enabled")


def notebook_extractor(ref: str) -> dict:
    try:
        import json as _json

        got = get_blobstore().get(ref)
        if got is None:
            return _err("notebook_extractor", "blob not found", source_id=ref)
        _, data = got
        nb = _json.loads(data.decode("utf-8", "ignore"))
        cells = []
        for i, cell in enumerate(nb.get("cells", [])[:40]):
            src = "".join(cell.get("source", []))
            if src.strip():
                cells.append({"id": f"{ref[:8]}-c{i}", "text": src[:2000], "ordinal": i})
        return _extractor_shape(ref, "application/x-ipynb+json", text_chunks=cells)
    except Exception as e:
        return _err("notebook_extractor", f"{type(e).__name__}: {e}", source_id=ref)


def table_extractor(ref: str) -> dict:
    return {"tables": [], "note": "structured table extraction not enabled"}


def code_repo_reader(url: str) -> dict:
    return _extractor_shape(url, "text/x-repo", note="repo reading not enabled")


def ocr_tool(ref: str) -> dict:
    return vision_image_extractor(ref)


TOOLS = {
    "web_search": web_search,
    "web_scrape": web_scrape,
    "paper_search": paper_search,
    "paper_downloader": paper_downloader,
    "file_fetcher": file_fetcher,
    "pdf_extractor": pdf_extractor,
    "slides_extractor": slides_extractor,
    "doc_extractor": doc_extractor,
    "html_article_extractor": html_article_extractor,
    "illustration_scraper": illustration_scraper,
    "vision_image_extractor": vision_image_extractor,
    "transcript_extractor": transcript_extractor,
    "notebook_extractor": notebook_extractor,
    "table_extractor": table_extractor,
    "code_repo_reader": code_repo_reader,
    "ocr_tool": ocr_tool,
}
