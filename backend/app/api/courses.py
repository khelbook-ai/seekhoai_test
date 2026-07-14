"""Course-build API (spec 07 §3-5). Drives the HITL build: create → clarify → cost
approval → (background) generation. The learning runtime is a separate router."""
from __future__ import annotations

import io
from pathlib import Path

from fastapi import APIRouter, BackgroundTasks, File, Form, HTTPException, UploadFile
from pydantic import BaseModel

from app import build, events
from app.db import fetchall
from app.guardrail import check
from app.store import get_clarifications, get_course, list_subtopics

router = APIRouter(prefix="/api/courses", tags=["courses"])


def _extract_material(filename: str, data: bytes) -> str:
    """Pull text out of an uploaded PDF / slide deck / doc so the Architect can build a
    course from it (spec 05 §12). Returns "" when nothing usable is found."""
    ext = Path(filename or "").suffix.lower()
    try:
        if ext == ".pdf":
            import pypdf
            r = pypdf.PdfReader(io.BytesIO(data))
            return "\n".join((p.extract_text() or "") for p in r.pages)
        if ext in (".pptx",):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        if ext in (".docx",):
            from docx import Document
            return "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)
        if ext in (".txt", ".md"):
            return data.decode("utf-8", "ignore")
    except Exception:
        return ""
    return ""


@router.get("")
def list_courses(user_id: str | None = None) -> dict:
    """A learner's courses with their current stage — powers the left sidebar (spec 07 §0).
    Scoped to the signed-in user when provided."""
    if user_id:
        rows = fetchall(
            "SELECT id, title, status, currency_mode, created_at FROM courses "
            "WHERE user_id = %s ORDER BY created_at DESC", (user_id,))
    else:
        rows = fetchall(
            "SELECT id, title, status, currency_mode, created_at FROM courses ORDER BY created_at DESC")
    return {"courses": [{"course_id": str(r["id"]), "title": r["title"], "status": r["status"],
                         "currency_mode": r["currency_mode"],
                         "created_at": r["created_at"].isoformat() if r["created_at"] else None}
                        for r in rows]}


@router.get("/{course_id}/events")
def build_events(course_id: str, after: int = 0) -> dict:
    """Live build log (spec 07 §5). Poll with the last-seen id in `after`."""
    return {"events": events.since(course_id, after_id=after)}


class CreateCourse(BaseModel):
    raw_prompt: str
    raw_role: str = ""
    user_id: str | None = None
    tester_override: bool = False
    override_reason: str | None = None


@router.post("")
def create_course(req: CreateCourse) -> dict:
    g = check(req.raw_prompt, "course_prompt", tester_override=req.tester_override,
              override_reason=req.override_reason)
    if not g.allow:
        raise HTTPException(400, g.user_message or "prompt rejected by guardrail")
    if req.raw_role:
        gr = check(req.raw_role, "clarify", tester_override=req.tester_override,
                   override_reason=req.override_reason)
        role = gr.sanitized_text if gr.allow else req.raw_role
    else:
        role = ""
    return build.start_build(raw_prompt=g.sanitized_text, raw_role=role, user_id=req.user_id)


@router.post("/from-file")
async def create_from_file(file: UploadFile = File(...), raw_role: str = Form(""),
                           user_id: str | None = Form(None)) -> dict:
    """Create a course from an uploaded PDF / slide deck / doc (spec 05 §12). The extracted
    text seeds the Architect, which builds the curriculum primarily from it."""
    data = await file.read()
    if len(data) > 20 * 1024 * 1024:
        raise HTTPException(400, "file too large (max 20MB)")
    material = _extract_material(file.filename or "", data)
    if len(material.strip()) < 200:
        raise HTTPException(400, "couldn't read enough text from that file — try a PDF, PPTX, DOCX, or text file")
    title = Path(file.filename or "uploaded material").stem.replace("-", " ").replace("_", " ").strip()
    prompt = f"Build a course from my uploaded material: \"{title}\""
    return build.start_build(raw_prompt=prompt, raw_role=raw_role or "",
                             user_id=user_id or None, seed_material=material[:200000])


class ClarifyReq(BaseModel):
    answers: dict[str, str]  # {ordinal_str: answer}


@router.post("/{course_id}/clarify")
def clarify(course_id: str, req: ClarifyReq) -> dict:
    # guard each free-text answer (neutralised, never blocked at this entry point)
    clean: dict[int, str] = {}
    for k, v in req.answers.items():
        g = check(v, "clarify")
        clean[int(k)] = g.sanitized_text
    try:
        return build.submit_clarifications(course_id, clean)
    except KeyError:
        raise HTTPException(404, "course not found")


@router.get("/{course_id}")
def get_course_detail(course_id: str) -> dict:
    course = get_course(course_id)
    if course is None:
        raise HTTPException(404, "course not found")
    clar = get_clarifications(course_id)
    role_row = fetchall("SELECT role_raw FROM users WHERE id = %s", (course["user_id"],)) if course["user_id"] else []
    return {
        "course_id": course_id,
        "title": course["title"],
        "status": course["status"],
        "currency_mode": course["currency_mode"],
        "raw_prompt": course["raw_prompt"],
        "raw_role": role_row[0]["role_raw"] if role_row else None,
        "clarifications": [{"ordinal": c["ordinal"], "q": c["question"],
                            "options": c["options"] or [], "answer": c["answer"],
                            "multi_select": c.get("multi_select", False)} for c in clar],
        "curriculum": course["curriculum"],
        "cost_estimate": course["cost_estimate"],
        "cost_approved": course["cost_approved"],
        "cost_actual": float(course["cost_actual"]) if course["cost_actual"] is not None else None,
        "cost_delta_abs": float(course["cost_delta_abs"]) if course["cost_delta_abs"] is not None else None,
        "cost_delta_pct": float(course["cost_delta_pct"]) if course["cost_delta_pct"] is not None else None,
        "cost_reconciliation": course["cost_reconciliation"],
        "subtopics": list_subtopics(course_id),
    }


class CostApproval(BaseModel):
    approved: bool


@router.post("/{course_id}/cost-approval")
def cost_approval(course_id: str, req: CostApproval, bg: BackgroundTasks) -> dict:
    try:
        result = build.approve_cost(course_id, req.approved)
    except KeyError:
        raise HTTPException(404, "course not found")
    if result["status"] == "building":
        bg.add_task(build.run_build, course_id)  # Phase 2 pipeline runs in background
    return result


class RestartReq(BaseModel):
    mode: str = "resume"  # "resume" last session (default) or "fresh" session over same content
    user_id: str | None = None


@router.post("/{course_id}/restart")
def restart_with_content(course_id: str, req: RestartReq) -> dict:
    """Token-free replay (spec 06 §6): relaunch a learning session bound to the already
    built course. NEVER re-invokes the build pipeline."""
    from app import runtime
    from app.db import fetchone

    course = get_course(course_id)
    if course is None:
        raise HTTPException(404, "course not found")
    if req.mode == "resume":
        last = fetchone(
            "SELECT id FROM sessions WHERE course_id = %s ORDER BY started_at DESC LIMIT 1",
            (course_id,))
        if last:
            sid = str(last["id"])
            return {"session_id": sid, "course_id": course_id, "resumed": True,
                    "interaction": runtime.current_interaction(sid)}
    return runtime.create_session(course_id, req.user_id)
